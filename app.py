
from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
presentation = Presentation()

# Slide 1: Title Slide
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sample Presentation"
subtitle.text = "An Example Using python-pptx"

# Slide 2: Title and Content (Bullet Points)
bullet_slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.shapes.placeholders[1]

title.text = "Key Features"
tf = body.text_frame
tf.text = "Create presentations programmatically."

p = tf.add_paragraph()
p.text = "Add slides with different layouts."
p.level = 1

p = tf.add_paragraph()
p.text = "Insert text, images, and tables."
p.level = 1

p = tf.add_paragraph()
p.text = "Customize text formatting and object positions."
p.level = 0

# Slide 3: Title Only (and add a text box manually)
title_only_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "Further Information"

left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "This presentation was generated entirely with Python using the python-pptx library."

# Slide 4: Title and Picture
pic_slide_layout = presentation.slide_layouts[8] # Assuming a layout with space for a picture
slide = presentation.slides.add_slide(pic_slide_layout)
title = slide.shapes.title
title.text = "A Placeholder Image"

# Add a placeholder image (you'd replace 'path/to/image.png' with a real image path)
# For this example, we'll skip adding a real image to avoid dependency on local files.
# If you have an image, uncomment the following lines and replace the path:
# img_path = "path/to/your/image.png"
# try:
#     left = Inches(1)
#     top = Inches(2.5)
#     pic = slide.shapes.add_picture(img_path, left, top, height=Inches(4.5))
# except FileNotFoundError:
#     # Add a text box instead if image not found
#     left = Inches(1)
#     top = Inches(2.5)
#     width = Inches(8)
#     height = Inches(1)
#     textbox = slide.shapes.add_textbox(left, top, width, height)
#     tf = textbox.text_frame
#     tf.text = "Image placeholder (replace 'path/to/your/image.png' with a real image path to see an image)."

# Add a text box to explain the image placeholder
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "No image included in this example. To add an image, replace this text box with `slide.shapes.add_picture('your_image.png', left, top)`"


# Slide 5: Title and Table
title_and_content_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(title_and_content_layout)
title = slide.shapes.title
title.text = "Sample Data Table"

rows = 3
cols = 3
left = Inches(1.5)
top = Inches(2.5)
width = Inches(7)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(2.5)

# Write table headers
table.cell(0, 0).text = "Header 1"
table.cell(0, 1).text = "Header 2"
table.cell(0, 2).text = "Header 3"

# Write data
table.cell(1, 0).text = "Row 1, Col 1"
table.cell(1, 1).text = "Row 1, Col 2"
table.cell(1, 2).text = "Row 1, Col 3"
table.cell(2, 0).text = "Row 2, Col 1"
table.cell(2, 1).text = "Row 2, Col 2"
table.cell(2, 2).text = "Row 2, Col 3"


# Save the presentation
presentation.save("output.pptx")
